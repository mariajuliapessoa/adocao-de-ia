import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

anos = ['2023', '2024']
percentual = [55, 72]
plt.figure(figsize=(6,4))
plt.bar(anos, percentual, color=['blue','green'])
plt.ylim(50,80)
plt.title('Adoção de IA pelas Empresas (%)')
plt.ylabel('Percentual de Empresas')
plt.savefig('grafico_adocao_ia.png')
plt.close()

# Vendas Online
meses = np.arange(1, 13)
vendas_sem_ia = [100, 105, 110, 115, 120, 125, 130, 128, 132, 135, 140, 145]
vendas_com_ia = [100, 108, 118, 130, 145, 160, 175, 180, 190, 200, 215, 230]
plt.figure(figsize=(8,5))
plt.bar(meses - 0.2, vendas_sem_ia, width=0.4, label='Sem IA', color='salmon')
plt.bar(meses + 0.2, vendas_com_ia, width=0.4, label='Com IA', color='skyblue')
plt.xticks(meses)
plt.xlabel('Mês')
plt.ylabel('Vendas (R$ mil)')
plt.title('Crescimento de Vendas Online com IA')
plt.legend()
plt.tight_layout()
plt.savefig('grafico_vendas_ia.png')
plt.close()

categorias = ['Antes IA', 'Depois IA']
estoque = [30, 12]
plt.figure(figsize=(6,4))
plt.bar(categorias, estoque, color=['red','green'])
plt.title('Otimização de Estoque com IA (%)')
plt.ylabel('% de Estoque Parado')
plt.savefig('grafico_estoque.png')
plt.close()

data_empresas = pd.read_csv('multiTimeline.csv', skiprows=1)
if 'isPartial' in data_empresas.columns:
    data_empresas = data_empresas.drop(columns=['isPartial'])
data_empresas[data_empresas.columns[0]] = pd.to_datetime(data_empresas[data_empresas.columns[0]], errors='coerce')
data_empresas = data_empresas.set_index(data_empresas.columns[0])
data_empresas = data_empresas.dropna()

plt.figure(figsize=(10,6))
for col in data_empresas.columns:
    plt.plot(data_empresas.index, data_empresas[col], label=col)
plt.title('Popularidade de Empresas com IA x Sem IA')
plt.ylabel('Interesse relativo')
plt.xlabel('Meses')
plt.legend()
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig('grafico_ia_vs_naoia.png')
plt.close()

prs = Presentation()

def add_textbox(slide, text, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4), font_size=Pt(18)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size
    return textbox

slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "IA Aplicada ao Varejo de Moda"
add_textbox(slide, "Protótipo de previsão de vendas, chatbot e cases de sucesso", top=Inches(3), font_size=Pt(20))

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Objetivos da IA"
add_textbox(slide,
    "• Aumentar vendas com recomendações personalizadas\n"
    "• Otimizar estoque e reduzir perdas\n"
    "• Melhorar atendimento ao cliente com chatbots\n"
    "• Analisar tendências de moda para decisões estratégicas"
)

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Benefícios Estratégicos"
add_textbox(slide,
    "• Decisões mais rápidas e precisas\n"
    "• Fidelização de clientes\n"
    "• Redução de custos operacionais\n"
    "• Diferencial competitivo no mercado"
)

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Cases de Sucesso"
add_textbox(slide,
    "• Zara: previsão de demanda e ajuste de produção\n"
    "• H&M: personalização de ofertas e experiência do cliente\n"
    "• Magazine Luiza: recomenda produtos online baseado no histórico\n"
    "• Renner: otimização de estoque e campanhas promocionais"
)

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Adoção de IA nas Empresas (72% em 2024)"
slide.shapes.add_picture('grafico_adocao_ia.png', Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Crescimento de Vendas Online com IA"
slide.shapes.add_picture('grafico_vendas_ia.png', Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Otimização de Estoque com IA"
slide.shapes.add_picture('grafico_estoque.png', Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Tendências de Moda (Google Trends Brasil)"
slide.shapes.add_picture('grafico_tendencias_reais.png', Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Popularidade de Empresas com IA x Sem IA"
slide.shapes.add_picture('grafico_ia_vs_naoia.png', Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Protótipo de Chatbot"
add_textbox(slide,
    "• Atende dúvidas sobre produtos, horários e entregas\n"
    "• Recomenda roupas personalizadas\n"
    "• Pode ser integrado a WhatsApp, Instagram ou e-commerce\n"
    "• Demonstra IA aplicada ao atendimento"
)

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Próximos Passos"
add_textbox(slide,
    "• Implementar chatbot e monitorar resultados\n"
    "• Criar modelos de previsão de vendas avançados\n"
    "• Integrar IA ao e-commerce e estoque\n"
    "• Avaliar ROI e expandir o uso de IA"
)

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Fontes"
add_textbox(slide,
    "1. CNN Brasil: Uso de IA atinge 72% das empresas (2024)\n"
    "2. InfoMoney: IA ajuda mais do que substitui pessoas, dizem executivos\n"
    "3. Gazeta do Povo: IA aumenta receita de empresas brasileiras\n"
    "4. FGV: Estudo sobre como IA pode aumentar desempenho de empresas\n"
    "5. Google Trends: Dados de interesse por categoria de produto e empresas\n"
    "6. McKinsey: Impacto da IA em marketing, vendas e cadeia de suprimentos"
)

prs.save("Apresentacao_IA_Varejo_Completo.pptx")
print("PPT completo gerado com todos os gráficos e fontes!")
